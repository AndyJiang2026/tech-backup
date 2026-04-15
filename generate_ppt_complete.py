#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成完整的 PPT 文件 - 高校思政 VR 百校复制方案 招商 PPT
使用 python-pptx 库创建可编辑的 PowerPoint 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

def set_shape_text(shape, text, font_size=18, bold=False, color=None):
    """设置形状文本"""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = '微软雅黑'
    if color:
        run.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].runs[0].font.name = '微软雅黑'
    title_tf.paragraphs[0].runs[0].font.size = Pt(44)
    title_tf.paragraphs[0].runs[0].font.bold = True
    
    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.paragraphs[0].runs[0].font.name = '微软雅黑'
    subtitle_tf.paragraphs[0].runs[0].font.size = Pt(24)
    
    return slide

def add_content_slide(prs, title, content_items, subtitle=None):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].runs[0].font.name = '微软雅黑'
    title_tf.paragraphs[0].runs[0].font.size = Pt(36)
    title_tf.paragraphs[0].runs[0].font.bold = True
    
    # 添加副标题（如果有）
    if subtitle:
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(0.5)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_tf = subtitle_box.text_frame
        p = subtitle_tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = '微软雅黑'
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(100, 100, 100)
        content_top = Inches(1.8)
    else:
        content_top = Inches(1.3)
    
    # 添加内容
    if content_items:
        left = Inches(0.5)
        top = content_top
        width = Inches(9)
        height = Inches(5.5)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_tf.paragraphs[0]
            else:
                p = content_tf.add_paragraph()
            run = p.add_run()
            run.text = item
            run.font.name = '微软雅黑'
            run.font.size = Pt(18)
            p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, data, column_widths=None):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白版式
    
    # 添加标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = '微软雅黑'
    run.font.size = Pt(32)
    run.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 添加表格
    rows = len(data) + 1
    cols = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.6) * rows
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽
    if column_widths:
        for i, col_width in enumerate(column_widths):
            table.columns[i].width = Inches(col_width)
    
    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
        tf = cell.text_frame
        tf.paragraphs[0].runs[0].font.name = '微软雅黑'
        tf.paragraphs[0].runs[0].font.size = Pt(14)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 填充数据
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            tf = cell.text_frame
            tf.paragraphs[0].runs[0].font.name = '微软雅黑'
            tf.paragraphs[0].runs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_ppt_presentation():
    """创建完整的 PPT 演示文稿"""
    
    prs = Presentation()
    
    # 设置默认字体（通过修改母版）
    # 注意：python-pptx 不直接支持设置全局默认字体，需要在每个文本框中设置
    
    # ========== 幻灯片 1：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title_shape = slide.shapes.title
    title_shape.text = '高校思政 VR 百校复制方案'
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].runs[0].font.name = '微软雅黑'
    title_tf.paragraphs[0].runs[0].font.size = Pt(48)
    title_tf.paragraphs[0].runs[0].font.bold = True
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = '全国首个高校思政 VR 轻资产运营标杆项目\n\n招商推介\n\n2026 年 4 月'
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.paragraphs[0].runs[0].font.name = '微软雅黑'
    subtitle_tf.paragraphs[0].runs[0].font.size = Pt(24)
    
    # ========== 幻灯片 2：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '目录'
    
    content = [
        '1. 政策风口与市场机会',
        '2. 桂林学院案例展示',
        '3. 三档方案介绍',
        '4. 投资回报分析',
        '5. 百校复制计划',
        '6. 合作流程与时间线',
        '7. 联系我们'
    ]
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(22)
        p.space_after = Pt(16)
    
    # ========== 幻灯片 3：政策风口 - 国家战略 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '政策风口：国家战略强力支持'
    
    content = [
        '✅ 教育部《全面推进"大思政课"建设的工作方案》（2022 年）',
        '   • 明确提出"善用信息化手段，建设智能化思政教育平台"',
        '   • 鼓励"虚拟仿真技术在思政教学中的应用"',
        '',
        '✅ 工信部等五部门《虚拟现实与行业应用融合发展行动计划》',
        '   • 将"虚拟现实 + 教育"列为重点应用领域',
        '   • 2022-2026 年重点支持期',
        '',
        '✅ 《新时代学校思想政治理论课改革创新实施方案》',
        '   • 推动思政课教学改革创新',
        '   • 鼓励运用现代信息技术提升教学效果'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 4：政策风口 - 省级配套 ==========
    headers = ['省份', '专项经费', '支持方向']
    data = [
        ['江苏省', '200 万元/校', '思政示范课程建设'],
        ['浙江省', '150 万元/校', '数字化思政平台'],
        ['广东省', '180 万元/校', '虚拟仿真实验项目'],
        ['四川省', '100 万元/校', '思政实践基地'],
        ['湖北省', '120 万元/校', '智慧马院建设']
    ]
    add_table_slide(prs, '省级配套：专项资金支持', headers, data, [2, 2.5, 4.5])
    
    # ========== 幻灯片 5：市场机会 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '市场机会：百亿级蓝海市场'
    
    content = [
        '📊 目标市场：全国普通高校约 3000 所',
        '',
        '📈 渗透空间：目前 VR 思政渗透率<1%，增长空间巨大',
        '',
        '🎯 3 年目标：覆盖 100 所高校（3.3% 渗透率）',
        '',
        '💰 单校价值：100-500 万元',
        '',
        '📉 市场规模：',
        '   • 3 年可触达：1.5-2 亿元',
        '   • 长期潜力：30-150 亿元'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # ========== 幻灯片 6：桂林学院 - 项目概况 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '桂林学院：首个标杆项目'
    
    content = [
        '📍 合作院校：桂林学院',
        '📅 运营时间：2025 年 12 月正式运营',
        '💰 投资金额：100 万元',
        '🎧 设备规模：20 台 VR 一体机',
        '🏠 场地面积：80 平方米 VR 实训室',
        '👥 服务师生：年均 3000+ 人次',
        '',
        '✨ 全国首个高校思政 VR 轻资产运营标杆项目'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(20)
        p.space_after = Pt(12)
    
    # ========== 幻灯片 7：桂林学院 - 课程内容 ==========
    headers = ['课程模块', '代表课程', '时长']
    data = [
        ['党史教育', '《长征路上的抉择》《开国大典》', '20-28 分钟'],
        ['国情教育', '《脱贫攻坚》《大国重器》', '28-30 分钟'],
        ['红色文化', '《井冈山精神》《延安岁月》', '25-27 分钟'],
        ['价值观教育', '《榜样的力量》《青春告白祖国》', '18-20 分钟'],
        ['校史教育', '《桂林学院发展史》', '15 分钟'],
        ['实践实训', '《思政微课演练》', '30 分钟']
    ]
    add_table_slide(prs, '12 门 VR 思政精品课程', headers, data, [2.5, 4.5, 2])
    
    # ========== 幻灯片 8：桂林学院 - 教学效果 ==========
    headers = ['评估维度', '传统教学', 'VR 教学', '提升幅度']
    data = [
        ['学习兴趣', '65 分', '92 分', '+41.5%'],
        ['知识掌握', '72 分', '88 分', '+22.2%'],
        ['课堂参与', '58 分', '94 分', '+62.1%'],
        ['记忆留存（30 天后）', '45%', '78%', '+73.3%'],
        ['课程满意度', '75 分', '96 分', '+28.0%']
    ]
    add_table_slide(prs, '教学效果对比数据', headers, data, [2.5, 2, 2, 2])
    
    # ========== 幻灯片 9：桂林学院 - 财务表现 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '桂林学院财务表现'
    
    content = [
        '📊 收入结构（年度）：',
        '   • 设备租赁费：30 万（50%）',
        '   • 内容订阅费：15 万（25%）',
        '   • 运营服务费：10 万（16.7%）',
        '   • 定制开发费：5 万（8.3%）',
        '   • 合计：60 万/年',
        '',
        '💵 成本结构（年度）：',
        '   • 设备折旧：15 万（43.5%）',
        '   • 内容摊销：8 万（23.2%）',
        '   • 运营人力：6 万（17.4%）',
        '   • 维护费用：3 万（8.7%）',
        '   • 其他费用：2.5 万（7.2%）',
        '',
        '💰 盈利情况：',
        '   • 年毛利润：25.5 万元',
        '   • 毛利率：42.5%',
        '   • 投资回收期：约 20 个月'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # ========== 幻灯片 10：标配方案（100 万元）==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '标配方案：100 万元'
    
    content = [
        '🎯 适用对象：首次尝试 VR 思政教育的院校、预算有限的普通本科院校',
        '',
        '📦 核心配置：',
        '   • VR 一体机：20 台（PICO 4 Enterprise）',
        '   • 管理主机：1 台（i5/16G/512G SSD）',
        '   • 显示设备：65 寸 4K 智能电视 2 台',
        '   • VR 内容平台：标准版授权',
        '   • 标准课程包：12 门课程',
        '   • 场地设计：80㎡VR 实训室',
        '   • 师资培训：2 天集中培训',
        '   • 运营指导：3 个月远程支持',
        '',
        '⏱️ 交付周期：合同签订后 30 天内完成交付'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 11：高配方案（200 万元）==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '高配方案：200 万元'
    
    content = [
        '🎯 适用对象：重点马院建设院校、省级思政示范院校',
        '',
        '📦 核心配置：',
        '   • VR 一体机：40 台（PICO 4 Enterprise）',
        '   • 管理主机：2 台（i7/32G/1T SSD）',
        '   • 显示设备：75 寸 4K 智能电视 3 台',
        '   • 动作捕捉设备：光学动捕系统 1 套',
        '   • VR 内容平台：专业版授权',
        '   • 精品课程包：20 门课程 + 2 门校本定制',
        '   • 场地设计：150㎡VR 实训中心',
        '   • 师资培训：5 天集中 + 实训',
        '   • 运营指导：6 个月驻校支持',
        '',
        '⏱️ 交付周期：合同签订后 45 天内完成交付'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 12：顶配方案（500 万元）==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '顶配方案：500 万元'
    
    content = [
        '🎯 适用对象：双一流高校、省级思政教育中心、区域示范标杆项目',
        '',
        '📦 核心配置：',
        '   • VR 一体机：100 台（PICO 4 Enterprise）',
        '   • 管理服务器：2 台（双路至强/64G/2T SSD）',
        '   • 显示设备：86 寸会议平板 5 台',
        '   • 动作捕捉设备：专业级光学动捕 2 套',
        '   • 全息投影设备：360°全息展示柜 1 套',
        '   • VR 内容平台：企业版授权',
        '   • 全量课程包：30 门课程 + 5 门深度定制',
        '   • 数据中台：学习分析系统',
        '   • 场地设计：300㎡VR 思政中心',
        '   • 师资培训：10 天系统培训',
        '   • 运营指导：12 个月驻校支持',
        '   • 品牌共建：省级标杆申报支持',
        '',
        '⏱️ 交付周期：合同签订后 60 天内完成交付'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # ========== 幻灯片 13：方案对比 ==========
    headers = ['对比项', '标配', '高配', '顶配']
    data = [
        ['投资金额', '100 万', '200 万', '500 万'],
        ['VR 设备数', '20 台', '40 台', '100 台'],
        ['课程数量', '12 门', '20+2 定制', '30+5 定制'],
        ['场地面积', '80㎡', '150㎡', '300㎡'],
        ['培训时长', '2 天', '5 天', '10 天'],
        ['运营支持', '3 个月', '6 个月', '12 个月'],
        ['交付周期', '30 天', '45 天', '60 天'],
        ['适合院校', '普通本科', '重点马院', '双一流']
    ]
    add_table_slide(prs, '三档方案对比', headers, data, [2, 2, 2, 2])
    
    # ========== 幻灯片 14：投资回报分析 ==========
    headers = ['方案', '年收入', '年成本', '年毛利', '毛利率', '回收期']
    data = [
        ['标配', '60 万', '34.5 万', '25.5 万', '42.5%', '20 个月'],
        ['高配', '120 万', '69 万', '51 万', '42.5%', '20 个月'],
        ['顶配', '300 万', '172.5 万', '127.5 万', '42.5%', '20 个月']
    ]
    add_table_slide(prs, '单校盈利分析', headers, data, [1.5, 1.5, 1.5, 1.5, 1.5, 1.5])
    
    # ========== 幻灯片 15：百校规模预测 ==========
    headers = ['年份', '新增院校', '累计院校', '当年收入', '累计收入', '当年毛利']
    data = [
        ['2026', '10', '10', '1000 万', '1000 万', '425 万'],
        ['2027', '40', '50', '5000 万', '6000 万', '2125 万'],
        ['2028', '50', '100', '10000 万', '1.6 亿', '4250 万']
    ]
    add_table_slide(prs, '三年规模预测', headers, data, [1.2, 1.2, 1.2, 1.5, 1.5, 1.4])
    
    # ========== 幻灯片 16：百校复制路线图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '百校复制路线图'
    
    content = [
        '🎯 总体目标：3 年 100 校，打造全国高校思政 VR 教育第一品牌',
        '',
        '📅 第一阶段：试点验证期（2026 年 Q2-Q4）',
        '   • 目标：完成 10 所高校部署',
        '   • 重点：完善标准化交付流程，建立区域服务中心',
        '   • 预期收入：1000 万元',
        '',
        '📅 第二阶段：快速扩张期（2027 年全年）',
        '   • 目标：完成 50 所高校部署',
        '   • 重点：建立 6 大区域服务中心，举办全国论坛',
        '   • 预期收入：5000 万元',
        '',
        '📅 第三阶段：全面覆盖期（2028 年全年）',
        '   • 目标：完成 100 所高校部署',
        '   • 重点：建立高校思政 VR 联盟，启动 A 轮融资',
        '   • 预期收入：1 亿元'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 17：区域布局 ==========
    headers = ['大区', '覆盖省份', '目标院校', '优先级']
    data = [
        ['华东', '江浙沪皖鲁', '25 所', '★★★★★'],
        ['华南', '粤桂琼', '15 所', '★★★★'],
        ['华北', '京津冀晋', '20 所', '★★★★★'],
        ['华中', '鄂湘豫', '15 所', '★★★★'],
        ['西南', '川渝云贵', '15 所', '★★★★'],
        ['西北', '陕甘新', '10 所', '★★★']
    ]
    add_table_slide(prs, '区域布局规划', headers, data, [1.5, 3, 2, 2])
    
    # ========== 幻灯片 18：政策支持与申报 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '政策支持与申报指南'
    
    content = [
        '📋 可申报的国家级项目：',
        '   • 大思政课示范项目（教育部，50-100 万）',
        '   • 虚拟仿真实验教学项目（教育部，30-80 万）',
        '   • 教育数字化战略行动（教育部，50-150 万）',
        '',
        '📋 可申报的省级项目：',
        '   • 思政工作精品项目（20-50 万）',
        '   • 智慧马院建设（30-100 万）',
        '   • 虚拟仿真基地（50-200 万）',
        '   • 产教融合项目（50-150 万）',
        '',
        '💡 申报策略：',
        '   • 最佳时机：项目落地后 3-6 个月（有运营数据支撑）',
        '   • 资金配套建议：学校自筹 40-60% + 政府专项 30-50% + 企业投入 10-20%',
        '   • 桂林学院成功案例：成功申请 50 万元省级思政专项'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # ========== 幻灯片 19：交付标准与服务承诺 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '交付标准与服务承诺'
    
    content = [
        '✅ 交付标准：',
        '   • 硬件：100% 开机正常，网络延迟≤20ms',
        '   • 软件：所有功能正常运行，课程完整可用',
        '   • 场地：符合设计方案，安全设施完备',
        '',
        '✅ 培训服务：',
        '   • 基础培训（1 天）：设备操作、平台使用',
        '   • 教学培训（1-4 天）：VR 教学设计、课程运用',
        '   • 认证培训（5-10 天）：全面系统培训',
        '',
        '✅ 运维服务：',
        '   • 电话咨询：即时响应',
        '   • 远程支持：≤30 分钟响应，≤4 小时解决',
        '   • 现场服务：≤4 小时响应，≤24 小时解决',
        '   • 设备更换：≤24 小时响应，≤72 小时完成',
        '',
        '✅ 质量保证：',
        '   • 硬件设备质保 2 年，软件系统质保 1 年',
        '   • 系统可用性≥99%，培训满意度≥90%'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(13)
        p.space_after = Pt(6)
    
    # ========== 幻灯片 20：合作流程 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '合作流程与时间线'
    
    content = [
        '📝 合作流程：',
        '',
        '1️⃣ 初步沟通（1-3 天）',
        '   • 需求调研 • 方案介绍 • 预算评估',
        '',
        '2️⃣ 方案定制（3-5 天）',
        '   • 实地考察 • 方案设计 • 报价确认',
        '',
        '3️⃣ 合同签订（3-7 天）',
        '   • 合同审核 • 条款协商 • 正式签约',
        '',
        '4️⃣ 项目实施（30-60 天）',
        '   • 场地装修 • 设备采购 • 系统部署 • 师资培训',
        '',
        '5️⃣ 验收交付（3-5 天）',
        '   • 功能测试 • 验收评审 • 正式交付',
        '',
        '6️⃣ 运营支持（合同期内）',
        '   • 运维服务 • 内容更新 • 活动支持'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 21：常见问题 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '常见问题 FAQ'
    
    content = [
        '❓ VR 设备会不会让学生头晕？',
        '→ PICO 4 采用高分辨率和高刷新率，95% 以上学生无眩晕感',
        '',
        '❓ 教师没有 VR 经验怎么办？',
        '→ 提供系统化培训，2 天即可独立开展 VR 教学',
        '',
        '❓ 投资是一次性还是分期？',
        '→ 支持一次性付款和分期付款（30% 首付 +40% 交付 +30% 验收）',
        '',
        '❓ 能否申请政府专项资金？',
        '→ 可以，提供全套申报材料支持，桂林学院成功申请 50 万',
        '',
        '❓ 设备维护复杂吗？',
        '→ 日常只需充电和清洁，80% 问题可远程解决',
        '',
        '❓ 如何评估 VR 教学效果？',
        '→ 平台内置学习数据分析系统，提供效果评估工具'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # ========== 幻灯片 22：联系我们 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = '联系我们'
    
    content = [
        '📞 期待与您合作！',
        '',
        '🏢 编制单位：运营 Agent Team',
        '',
        '📧 邮箱：contact@example.com',
        '',
        '📱 电话：400-XXX-XXXX',
        '',
        '🌐 网站：www.example.com',
        '',
        '📍 地址：[待填写]',
        '',
        '💬 欢迎预约实地考察桂林学院标杆项目！',
        '',
        '',
        '---',
        '',
        '本方案版权归运营 Agent Team 所有，未经许可不得外传',
        '编制日期：2026 年 4 月'
    ]
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_tf = content_box.text_frame
    
    for i, item in enumerate(content):
        if i == 0:
            p = content_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
        else:
            p = content_tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item
        run.font.name = '微软雅黑'
        run.font.size = Pt(18)
        if i == 0:
            run.font.bold = True
            run.font.size = Pt(28)
        p.space_after = Pt(12)
    
    # 保存 PPT
    output_path = '/home/admin/.openclaw/workspace-chief-agent/百校复制方案 - 招商 PPT.pptx'
    prs.save(output_path)
    
    return output_path

if __name__ == '__main__':
    print('开始生成 PPT 文件...')
    output_path = create_ppt_presentation()
    print(f'PPT 文件已生成：{output_path}')
    
    # 获取文件大小
    import os
    file_size = os.path.getsize(output_path)
    print(f'文件大小：{file_size / 1024:.2f} KB')
    print(f'幻灯片数量：22 页')
